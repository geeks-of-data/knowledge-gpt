import pptx
import pandas as pd


def process_pptx(pptx_file):
    """
    Extracts paragraphs from a Powerpoint file
    :param pptx_file: Powerpoint file
    :return: Dataframe with paragraphs
    """

    prs = pptx.Presentation(pptx_file)
    
    df_cook = pd.DataFrame(columns=['slide_number', 'content'])

    for slide_num, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        text = "\n".join(text)
        df_cook = df_cook.append({'slide_number': slide_num+1, 'content': text}, ignore_index=True)

    return df_cook
